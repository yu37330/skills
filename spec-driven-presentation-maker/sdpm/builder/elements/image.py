# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Image element."""
import sys
from pathlib import Path
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from sdpm.schema.defaults import ELEMENT_DEFAULTS
from sdpm.utils.image import resolve_image_path, apply_image_effects
from sdpm.utils.effects import apply_effects
from sdpm.utils.svg import _recolor_svg, get_svg_dimensions, generate_qr_svg, add_svg_to_slide
from sdpm.utils.text import _expand_styled_newlines, parse_styled_text
from sdpm.assets import is_recolor_protected

_DEFAULTS = ELEMENT_DEFAULTS["image"]


def _plain_label_text(line: str) -> str:
    """Return a label line's visible text with any {{...:text}} style markup
    stripped, for width estimation."""
    try:
        return "".join(seg.get("text", "") for seg in parse_styled_text(line))
    except Exception:
        return line


def _label_line_em(line: str) -> float:
    """Return a label line's estimated display width in em at the label font
    size. Latin/half-width glyphs count as 0.85em (calibrated against
    PowerPoint's rendering, see the call site); East Asian wide/fullwidth
    glyphs (CJK) are square, so they count as 1.05em — a plain character
    count would under-size Japanese labels and re-introduce the one-glyph-
    per-line wrapping this box is sized to prevent. Style markup
    ({{...:text}}) is stripped first."""
    import unicodedata

    text = _plain_label_text(line)
    return sum(1.05 if unicodedata.east_asian_width(ch) in ("W", "F") else 0.85
               for ch in text)

class ImageMixin:
    """Mixin providing image element methods."""

    def _add_image(self, slide, elem):
        """Add image element to slide.
        
        src: icons:NAME or file path (supports ~)
        """
        from pptx.enum.text import PP_ALIGN
        
        src = elem.get("src") or elem.get("path", "")
        x_pct = elem.get("x", 0)
        y_pct = elem.get("y", 0)
        width_pct = elem.get("width")
        height_pct = elem.get("height")
        label = elem.get("label")
        label_pos = elem.get("labelPosition", "bottom")
        label_size = elem.get("labelSize", 11)
        link = elem.get("link")
        rotation = elem.get("rotation", _DEFAULTS["rotation"])
        icon_color = elem.get("iconColor")

        if not src:
            print("Warning: image element without 'src' skipped — element dropped from PPTX", file=sys.stderr)
            return
        
        # QR code generation
        if src.startswith("qr:"):
            qr_url = src[3:]
            qr_size = int(width_pct or height_pct or 200)
            svg_bytes = generate_qr_svg(
                qr_url, size=qr_size,
                color=elem.get("color"),
                gradient=elem.get("gradient"),
                theme="dark" if self.is_dark else "light",
            )
            x = self._px_to_emu(x_pct)
            y = self._px_to_emu(y_pct)
            w = self._px_to_emu(width_pct or 200)
            h = self._px_to_emu(height_pct or 200)
            pic = add_svg_to_slide(slide, svg_bytes, x, y, w, h)
            if rotation != 0 and pic is not None:
                xfrm = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rotation * 60000)))
            # Label for QR (reuse inline label logic)
            if label and label_pos != "none":
                from pptx.enum.text import PP_ALIGN
                label_margin = int(h * 0.04)
                if label_pos == "bottom":
                    lbl_x, lbl_y, lbl_w = x, y + h + label_margin, w
                elif label_pos == "right":
                    lbl_x, lbl_y, lbl_w = x + w + Emu(100000), y + h // 3, self._px_to_emu(15)
                else:
                    return
                textbox = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, Emu(300000))
                tf = textbox.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
                self._apply_styled_text(p, label, default_font_size=label_size)
            return
        
        # Resolve image path
        if src.startswith("icons:") or src.startswith("assets:"):
            img_path = resolve_image_path(src, "dark" if self.is_dark else "light")
        else:
            img_path = Path(src).expanduser()
            if not img_path.is_absolute():
                img_path = self._base_dir / src
            if not img_path.exists():
                print(f"Warning: Image not found: {img_path}", file=sys.stderr)
                return
        
        is_svg = img_path.suffix.lower() == ".svg"

        # Prepare SVG bytes (with optional recolor)
        svg_bytes = None
        if is_svg:
            svg_bytes = img_path.read_bytes()
            if icon_color == "none":
                # Explicit opt-out: keep the SVG's own colors (used for
                # artwork imported from existing decks, not theme icons).
                pass
            elif src and is_recolor_protected(src):
                if icon_color:
                    print(f"Warning: iconColor ignored (recolor-protected asset): {src}", file=sys.stderr)
            else:
                effective_icon_color = icon_color or self.theme_colors["text"]
                recolored = _recolor_svg(svg_bytes, effective_icon_color)
                if recolored:
                    svg_bytes = recolored
        elif icon_color:
            print(f"Warning: iconColor ignored (not SVG): {img_path.name}", file=sys.stderr)
        
        # Calculate dimensions
        x = self._px_to_emu(x_pct)
        y = self._px_to_emu(y_pct)
        fit = elem.get("fit", _DEFAULTS["fit"])
        cover_crop = None
        
        if width_pct:
            width = self._px_to_emu(width_pct)
            if height_pct:
                height = self._px_to_emu(height_pct)
                # Apply fit logic when both dimensions specified
                if is_svg:
                    img_w, img_h = get_svg_dimensions(img_path)
                else:
                    from PIL import Image
                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                    except Exception:
                        img_w, img_h = 1, 1
                if img_w > 0 and img_h > 0 and not elem.get("crop"):
                    # With an explicit crop, srcRect fills the frame
                    # (PowerPoint semantics) — no contain/cover adjustment.
                    img_ratio = img_w / img_h
                    box_ratio = width / height
                    # Warn if aspect ratios differ significantly
                    if abs(img_ratio - box_ratio) / max(img_ratio, box_ratio) > 0.05:
                        sw_h = int(width_pct / img_ratio)  # width-based height
                        sh_w = int(height_pct * img_ratio)  # height-based width
                        if fit == "cover":
                            consequence = "image is being cropped"
                        elif fit == "stretch":
                            consequence = "image is being distorted"
                        else:
                            consequence = "image has padding within the box"
                        print(f"Warning: {consequence} (fit={fit}): {src}\n"
                              f"  box {width_pct}×{height_pct} → to preserve aspect ratio: "
                              f"width={width_pct}, height={sw_h} or "
                              f"width={sh_w}, height={height_pct}",
                              file=sys.stderr)
                    if fit == "contain":
                        if img_ratio > box_ratio:
                            height = int(width / img_ratio)
                        else:
                            width = int(height * img_ratio)
                    elif fit == "cover":
                        if img_ratio > box_ratio:
                            crop_pct = (1 - box_ratio / img_ratio) / 2
                            cover_crop = {"l": crop_pct, "r": crop_pct, "t": 0, "b": 0}
                        else:
                            crop_pct = (1 - img_ratio / box_ratio) / 2
                            cover_crop = {"l": 0, "r": 0, "t": crop_pct, "b": crop_pct}
                        # Keep width/height as the box size
                        width = self._px_to_emu(width_pct)
                        height = self._px_to_emu(height_pct)
            else:
                # Maintain original aspect ratio
                if is_svg:
                    img_w, img_h = get_svg_dimensions(img_path)
                else:
                    from PIL import Image
                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                    except Exception:
                        img_w, img_h = 1, 1
                height = int(width * img_h / img_w) if img_w > 0 else width
        
        if is_svg:
            if not width_pct:
                img_w, img_h = get_svg_dimensions(img_path)
                width = self._px_to_emu(img_w)
                height = self._px_to_emu(img_h)
            pic = add_svg_to_slide(slide, svg_bytes, x, y, width, height)
        else:
            if width_pct:
                pic = slide.shapes.add_picture(str(img_path), x, y, width=width, height=height)
            else:
                pic = slide.shapes.add_picture(str(img_path), x, y)
                width = pic.width
                height = pic.height
        
        # Apply cover crop via srcRect
        if cover_crop and pic is not None:
            from lxml import etree
            from pptx.oxml.ns import qn
            pic_el = pic._element if hasattr(pic, '_element') else pic
            blip_fill = pic_el.find(qn('p:blipFill'))
            if blip_fill is None:
                blip_fill = pic_el.find(qn('pic:blipFill'))
            if blip_fill is not None:
                for sr in blip_fill.findall(qn('a:srcRect')):
                    blip_fill.remove(sr)
                src_rect = etree.Element(qn('a:srcRect'))
                for attr, key in [('l', 'l'), ('t', 't'), ('r', 'r'), ('b', 'b')]:
                    if cover_crop[key]:
                        src_rect.set(attr, str(int(cover_crop[key] * 100000)))
                blip = blip_fill.find(qn('a:blip'))
                if blip is not None:
                    blip.addnext(src_rect)
                else:
                    blip_fill.insert(0, src_rect)
        
        # Apply line (border) — raster pictures only
        if not is_svg and pic is not None:
            line_color = elem.get("line")
            if line_color and line_color != "none":
                pic.line.fill.solid()
                hex_color = line_color.lstrip("#")
                pic.line.color.rgb = RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16),
                )
                pic.line.width = Pt(elem.get("lineWidth", 1))

        # Apply flip (mirrored pictures — cutout photos rely on this to face
        # the right way). pic may be a Picture or a raw XML element (SVG path).
        if (elem.get("flipH") or elem.get("flipV")) and pic is not None:
            pic_el = pic._element if hasattr(pic, '_element') else pic
            xfrm = pic_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                if elem.get("flipH"):
                    xfrm.set('flipH', '1')
                if elem.get("flipV"):
                    xfrm.set('flipV', '1')

        # Apply rotation
        if rotation != 0:
            if is_svg:
                # Set rotation directly on XML xfrm element
                xfrm = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rotation * 60000)))
            else:
                pic.rotation = rotation
        
        # Add hyperlink if specified
        if link:
            if is_svg:
                pass  # Hyperlinks on SVG not supported via direct XML (rare use case)
            else:
                pic.click_action.hyperlink.address = link
        
        # Add label if specified
        if label and label_pos != "none":
            # Scale margin proportionally to icon size (base: 4% of height)
            label_margin = int(height * 0.04)
            if label_pos == "bottom":
                # The label box must be WIDER than the icon, or a caption like
                # "Cognito" wraps one glyph per line ("Co / gni / to") inside a
                # 60px icon-width box even with word_wrap off. Size the box to the
                # longest label line and center it on the icon, so the caption
                # stays a single readable line that overhangs the icon evenly.
                _lbl_lines = _expand_styled_newlines(label.replace("\\n", "\n")).split("\n")
                _max_em = max(
                    (_label_line_em(ln) for ln in _lbl_lines), default=0.0)
                # Width per glyph at the label font size. PowerPoint renders
                # wider than a naive em estimate (kerning + internal box
                # margins), so _label_line_em uses ~0.85em for Latin and
                # ~1.05em for CJK fullwidth glyphs, plus generous padding —
                # a box even slightly too narrow wraps mid-word. Overhang past
                # the icon is harmless (labels are centered and the engine
                # leaves margin); a too-narrow box is not.
                _text_w_px = int(_max_em * label_size) + 16
                _icon_w_px = width_pct or 60
                lbl_w_px = max(_icon_w_px, _text_w_px)
                lbl_w = self._px_to_emu(lbl_w_px)
                lbl_x = x + width // 2 - lbl_w // 2
                lbl_y = y + height + label_margin
            elif label_pos == "right":
                lbl_x = x + width + label_margin * 2
                lbl_y = y + height // 3
                lbl_w = self._px_to_emu(15)
            else:
                return
            
            textbox = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, Emu(300000))
            tf = textbox.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
            label = label.replace("\\n", "\n")
            label = _expand_styled_newlines(label)
            lines = label.split("\n")
            align = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
            for li, line in enumerate(lines):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = align
                self._apply_styled_text(p, line, default_font_size=label_size)
        
        # Apply visual effects to image
        if pic is not None:
            pic_el = pic._element if hasattr(pic, '_element') else pic
            apply_image_effects(pic_el, elem)
            apply_effects(pic_el, elem, self.EMU_PER_PX)
    

    def _add_video(self, slide, elem):
        """Add video element to slide."""
        src = elem.get("src", "")
        poster = elem.get("poster", "")
        if not src:
            print("Warning: video element without 'src' skipped — element dropped from PPTX", file=sys.stderr)
            return
        
        video_path = self._base_dir / src
        poster_path = self._base_dir / poster if poster else None
        
        if not video_path or not Path(video_path).exists():
            print(f"Warning: Video not found: {src}", file=sys.stderr)
            return
        
        x = Emu(elem.get("_xEmu") or self._px_to_emu(elem.get("x", 0)))
        y = Emu(elem.get("_yEmu") or self._px_to_emu(elem.get("y", 0)))
        w = Emu(elem.get("_widthEmu") or self._px_to_emu(elem.get("width", 100)))
        h = Emu(elem.get("_heightEmu") or self._px_to_emu(elem.get("height", 100)))
        
        # Determine MIME type
        ext = Path(video_path).suffix.lower()
        mime_map = {'.mp4': 'video/mp4', '.avi': 'video/avi', '.wmv': 'video/x-ms-wmv', '.mov': 'video/quicktime'}
        mime = mime_map.get(ext, 'video/mp4')
        
        slide.shapes.add_movie(
            str(video_path),
            x, y, w, h,
            poster_frame_image=str(poster_path) if poster_path and Path(poster_path).exists() else None,
            mime_type=mime,
        )
