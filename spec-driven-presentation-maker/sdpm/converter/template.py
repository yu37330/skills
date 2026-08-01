# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Extract a placeholder-only template PPTX from an arbitrary input PPTX.

The output PPTX preserves every slide_master, slide_layout, and theme
from the source — nothing is removed from those parts. The original
slides are dropped and replaced with one placeholder-only sample slide
per layout that the source PPTX actually exercised. That way the
template stays browsable in PowerPoint (one preview per layout used),
remains a drop-in replacement for any layout the source used, *and*
keeps every other layout the deck author may want to reach for later.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation


_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# Extension URI of the PowerPoint section list (p14:sectionLst) inside
# presentation.xml's <p:extLst>.
_SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def _drop_section_list(presentation_el) -> None:
    """Remove the section-list ext from a <p:presentation> element.

    Args:
        presentation_el: The lxml ``<p:presentation>`` element
            (parent of ``<p:sldIdLst>``).
    """
    from pptx.oxml.ns import qn

    ext_lst = presentation_el.find(qn("p:extLst"))
    if ext_lst is None:
        return
    for ext in list(ext_lst.findall(qn("p:ext"))):
        if ext.get("uri") == _SECTION_EXT_URI:
            ext_lst.remove(ext)
    if len(ext_lst) == 0:
        presentation_el.remove(ext_lst)


def extract_placeholder_template(
    pptx_path: str | Path,
    output_path: str | Path,
) -> dict:
    """Build a placeholder-only template from a PPTX.

    Args:
        pptx_path: Source PPTX file.
        output_path: Destination path for the placeholder template.

    Returns:
        Metadata dict: {input_size, output_size, layout_count,
        master_count, used_layout_count}.

    Behavior:
        - All slide_masters are retained.
        - All slide_layouts are retained (including layouts the source
          PPTX never used).
        - All slides from the source are dropped, then one
          placeholder-only sample slide is emitted for each layout that
          at least one source slide used. Layouts the source never
          touched do not get sample slides.
    """
    src = Path(pptx_path)
    dst = Path(output_path)
    dst.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(src))

    # Step 1: figure out which layout parts are referenced by source slides.
    used_layout_partnames: set[str] = set()
    for slide in prs.slides:
        used_layout_partnames.add(str(slide.slide_layout.part.partname))

    # Step 2: drop every original slide (rel + sldIdLst entry).
    sldIdLst = prs.slides._sldIdLst
    for entry in list(sldIdLst):
        rId = entry.attrib.get(f"{_REL_NS}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        sldIdLst.remove(entry)

    # Step 2b: drop the section list (p14:sectionLst) if present. Sections
    # reference the original slide IDs; keeping them after the slides are
    # dropped leaves stale IDs that trigger PowerPoint's repair prompt.
    # Removing the whole ext (rather than patching individual entries) is
    # the safest option — section grouping is meaningless for a
    # placeholder-only template anyway.
    _drop_section_list(sldIdLst.getparent())

    # Step 3: emit one placeholder-only sample slide per *used* layout.
    # Iterate masters/layouts in their original order so the output
    # presentation order is stable.
    used_count = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if str(layout.part.partname) in used_layout_partnames:
                prs.slides.add_slide(layout)
                used_count += 1

    master_count = len(prs.slide_masters)
    layout_count = sum(len(m.slide_layouts) for m in prs.slide_masters)

    prs.save(str(dst))

    return {
        "input_size": src.stat().st_size,
        "output_size": dst.stat().st_size,
        "layout_count": layout_count,
        "master_count": master_count,
        "used_layout_count": used_count,
    }
